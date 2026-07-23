from __future__ import annotations
import os
import io
import uuid
import copy
import json
import math
import threading
import traceback
import re
import time
from pathlib import Path
import requests

from flask import Flask, request, jsonify, send_file, render_template
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree
import anthropic

# ---------------------------------------------------------------------------
# App setup
# ---------------------------------------------------------------------------

app = Flask(__name__)

BASE_DIR = Path(__file__).parent
UPLOAD_FOLDER = BASE_DIR / "uploads"
OUTPUT_FOLDER = BASE_DIR / "outputs"
UPLOAD_FOLDER.mkdir(exist_ok=True)
OUTPUT_FOLDER.mkdir(exist_ok=True)

# In-memory job registry
# {job_id: {"status": ..., "message": ..., "progress": ..., "plan": ..., "pptx_path": ..., "output": ...}}
jobs: dict = {}
jobs_lock = threading.Lock()


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------

def clone_slide(prs: Presentation, source_index: int = 0):
    """Clone the slide at source_index and append a copy to the presentation."""
    source = prs.slides[source_index]
    new_slide = prs.slides.add_slide(source.slide_layout)

    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)

    for child in source.shapes._spTree:
        sp_tree.append(copy.deepcopy(child))

    for rel in source.part.rels.values():
        if "image" in rel.reltype:
            try:
                new_slide.part.relate_to(rel.target_part, rel.reltype)
            except Exception:
                pass

    return new_slide


def replace_text_in_slide(slide, replacements: dict, ordered: dict = None):
    """
    Replace placeholder tokens in every text frame on a slide.

    replacements  – {old: new} for unique tokens
    ordered       – {old: [val1, val2, val3]} for tokens that appear
                    multiple times; replaced in document order (top→bottom)
    """
    order_counts = {k: 0 for k in (ordered or {})}

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in para.runs)
            if not full_text.strip():
                continue

            modified = full_text
            changed = False

            for placeholder, value in replacements.items():
                if placeholder in modified:
                    modified = modified.replace(
                        placeholder, str(value) if value is not None else ""
                    )
                    changed = True

            for placeholder, values in (ordered or {}).items():
                if placeholder in modified:
                    idx = order_counts[placeholder]
                    if idx < len(values):
                        modified = modified.replace(
                            placeholder,
                            str(values[idx]) if values[idx] is not None else "",
                        )
                        order_counts[placeholder] += 1
                        changed = True

            if changed and para.runs:
                para.runs[0].text = modified
                for run in para.runs[1:]:
                    run.text = ""


# ---------------------------------------------------------------------------
# Real landmark lookup — Google Maps (preferred) with OSM fallback
# ---------------------------------------------------------------------------

OSM_HEADERS = {"User-Agent": "Skyscale-OOH-Generator/1.0 (contact@skyscale.com)"}

def _haversine_km(lat1: float, lon1: float, lat2: float, lon2: float) -> float:
    R = 6371.0
    dlat = math.radians(lat2 - lat1)
    dlon = math.radians(lon2 - lon1)
    a = (math.sin(dlat / 2) ** 2
         + math.cos(math.radians(lat1)) * math.cos(math.radians(lat2))
         * math.sin(dlon / 2) ** 2)
    return R * 2 * math.asin(math.sqrt(a))


def _get_landmarks_google(location: str, city: str, n: int = 3):
    api_key = os.environ.get("GOOGLE_MAPS_API_KEY", "")
    if not api_key:
        return None
    try:
        geo = requests.get(
            "https://maps.googleapis.com/maps/api/geocode/json",
            params={"address": f"{location}, {city}", "key": api_key},
            timeout=8,
        ).json()
        if geo.get("status") != "OK":
            return None
        loc = geo["results"][0]["geometry"]["location"]
        lat, lng = loc["lat"], loc["lng"]

        places = requests.get(
            "https://maps.googleapis.com/maps/api/place/nearbysearch/json",
            params={
                "location": f"{lat},{lng}",
                "radius": 5000,
                "type": "point_of_interest",
                "key": api_key,
            },
            timeout=10,
        ).json()

        results = []
        seen: set = set()
        for p in places.get("results", []):
            name = p.get("name", "").strip()
            if not name or name in seen:
                continue
            seen.add(name)
            p_lat = p["geometry"]["location"]["lat"]
            p_lng = p["geometry"]["location"]["lng"]
            dist = _haversine_km(lat, lng, p_lat, p_lng)
            results.append((dist, name))

        results.sort(key=lambda x: x[0])
        filtered = [(d, name) for d, name in results if d <= 5.0][:n]
        return [
            f"{name} – {round(d, 1) if d >= 0.1 else 0.1}km"
            for d, name in filtered
        ] or None

    except Exception:
        return None


def _get_landmarks_osm(location: str, city: str, n: int = 3):
    try:
        geo_resp = requests.get(
            "https://nominatim.openstreetmap.org/search",
            params={"q": f"{location}, {city}", "format": "json", "limit": 1},
            headers=OSM_HEADERS,
            timeout=8,
        )
        geo_data = geo_resp.json()
        if not geo_data:
            return None
        lat = float(geo_data[0]["lat"])
        lon = float(geo_data[0]["lon"])
        time.sleep(1.1)

        overpass_query = f"""
[out:json][timeout:12];
(
  node["name"]["tourism"](around:5000,{lat},{lon});
  node["name"]["amenity"~"^(restaurant|cafe|hotel|bank|museum|theatre|cinema|hospital|university|library|historic)$"](around:5000,{lat},{lon});
  node["name"]["historic"](around:5000,{lat},{lon});
  node["name"]["shop"~"^(mall|department_store|supermarket)$"](around:5000,{lat},{lon});
);
out center 20;
"""
        elements = requests.post(
            "https://overpass-api.de/api/interpreter",
            data={"data": overpass_query},
            headers=OSM_HEADERS,
            timeout=15,
        ).json().get("elements", [])

        seen: set = set()
        ranked: list = []
        for el in elements:
            name = el.get("tags", {}).get("name", "").strip()
            if not name or name in seen:
                continue
            seen.add(name)
            el_lat = el.get("lat") or el.get("center", {}).get("lat", lat)
            el_lon = el.get("lon") or el.get("center", {}).get("lon", lon)
            ranked.append((_haversine_km(lat, lon, float(el_lat), float(el_lon)), name))

        ranked.sort(key=lambda x: x[0])
        results = []
        for dist, name in ranked:
            if dist > 5.0:
                break
            km = round(dist, 1) if dist >= 0.1 else 0.1
            results.append(f"{name} – {km}km")
            if len(results) == n:
                break

        return results if len(results) >= n else None

    except Exception:
        return None


def get_real_landmarks(location: str, city: str, n: int = 3):
    result = _get_landmarks_google(location, city, n)
    if result and len(result) >= n:
        return result
    return _get_landmarks_osm(location, city, n)


# ---------------------------------------------------------------------------
# Google Maps Static screenshot
# Adjust these constants to match your template layout (inches from top-left)
# ---------------------------------------------------------------------------

MAP_IMG_LEFT   = Inches(6.60)   # horizontal offset from slide left edge
MAP_IMG_TOP    = Inches(1.20)   # vertical offset from slide top edge
MAP_IMG_WIDTH  = Inches(3.10)   # image width
MAP_IMG_HEIGHT = Inches(2.10)   # image height


def get_map_image_bytes(location: str, city: str, zoom: int = 16) -> bytes | None:
    """
    Return a map image for the location.
    Uses Google Maps Static API when GOOGLE_MAPS_API_KEY is set,
    otherwise geocodes via Nominatim and fetches from staticmap.openstreetmap.de
    — no extra dependencies needed.
    """
    UA = "OOHProposalGenerator/1.0 (skyscalemedia.com)"
    api_key = os.environ.get("GOOGLE_MAPS_API_KEY", "")

    if api_key:
        try:
            geo = requests.get(
                "https://maps.googleapis.com/maps/api/geocode/json",
                params={"address": f"{location}, {city}", "key": api_key},
                timeout=8,
            ).json()
            if geo.get("status") != "OK" or not geo.get("results"):
                return None
            loc = geo["results"][0]["geometry"]["location"]
            lat, lng = loc["lat"], loc["lng"]
            resp = requests.get(
                "https://maps.googleapis.com/maps/api/staticmap",
                params={
                    "center": f"{lat},{lng}",
                    "zoom": zoom,
                    "size": "600x400",
                    "scale": 2,
                    "maptype": "roadmap",
                    "markers": f"color:red|size:mid|{lat},{lng}",
                    "key": api_key,
                },
                timeout=12,
            )
            ct = resp.headers.get("content-type", "")
            if resp.status_code == 200 and ct.startswith("image"):
                return resp.content
        except Exception:
            pass
        return None

    # ── Free fallback: Nominatim geocode → stitch OSM tiles with Pillow ──
    # Respect Nominatim's 1 req/sec rate limit
    time.sleep(1.1)
    try:
        geo = requests.get(
            "https://nominatim.openstreetmap.org/search",
            params={"q": f"{location}, {city}", "format": "json", "limit": 1},
            headers={"User-Agent": UA},
            timeout=8,
        ).json()
        if not geo:
            print(f"[MAP] Nominatim found nothing for {location!r}, {city!r}")
            return None
        lat, lng = float(geo[0]["lat"]), float(geo[0]["lon"])
        print(f"[MAP] geocoded {location!r} → lat={lat:.4f} lng={lng:.4f}")
    except Exception as e:
        print(f"[MAP] Nominatim failed for {location!r}: {e}")
        return None

    # Stitch a 3×3 grid of OSM tiles into one image
    try:
        from PIL import Image, ImageDraw

        z = min(zoom, 15)
        n = 2 ** z
        lat_rad = math.radians(lat)
        # Fractional tile position of the exact location
        fx = (lng + 180) / 360 * n
        fy = (1 - math.log(math.tan(lat_rad) + 1 / math.cos(lat_rad)) / math.pi) / 2 * n
        tx = int(fx)
        ty = int(fy)

        TILE = 256
        canvas = Image.new("RGB", (TILE * 3, TILE * 3), (220, 220, 215))
        ok = 0
        for row in range(3):
            for col in range(3):
                url = f"https://tile.openstreetmap.org/{z}/{tx - 1 + col}/{ty - 1 + row}.png"
                try:
                    r = requests.get(url, headers={"User-Agent": UA}, timeout=8)
                    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
                        tile = Image.open(io.BytesIO(r.content)).convert("RGB")
                        canvas.paste(tile, (col * TILE, row * TILE))
                        ok += 1
                except Exception:
                    pass

        print(f"[MAP] stitched {ok}/9 tiles for {location!r}")

        # Pin at the exact sub-pixel location within the 3×3 canvas
        pr = 13
        cx = int((fx - (tx - 1)) * TILE)
        cy = int((fy - (ty - 1)) * TILE)
        cx = max(pr + 3, min(canvas.width  - pr - 3, cx))
        cy = max(pr + 3, min(canvas.height - pr - 3, cy))

        draw = ImageDraw.Draw(canvas)
        # Soft shadow (solid dark grey offset)
        draw.ellipse([cx - pr + 3, cy - pr + 3, cx + pr + 3, cy + pr + 3],
                     fill=(100, 100, 100))
        # Red pin with white border
        draw.ellipse([cx - pr, cy - pr, cx + pr, cy + pr],
                     fill=(220, 30, 30), outline="white", width=3)

        buf = io.BytesIO()
        canvas.save(buf, format="PNG")
        return buf.getvalue()
    except Exception as e:
        print(f"[MAP] tile stitch failed for {location!r}: {e}")
        return None


def _cleanup_map_images(job_id: str):
    """Delete all map image files for a job."""
    with jobs_lock:
        map_paths = jobs.get(job_id, {}).get("map_paths", {})
    for path_str in map_paths.values():
        try:
            Path(path_str).unlink(missing_ok=True)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# AI content generation
# ---------------------------------------------------------------------------

def generate_site_content(site: dict, client: anthropic.Anthropic) -> dict:
    site_name = site.get("Site Name", "")
    location  = site.get("Location", "")
    market    = site.get("Market", "")
    fmt       = site.get("Format", "")
    size      = site.get("Size", "")
    is_mobile = str(location).strip().lower() == "various"

    lookup_address = market if is_mobile else location
    real_landmarks: list | None = get_real_landmarks(lookup_address, market)

    if real_landmarks:
        landmark_instruction = (
            "Real nearby landmarks have already been sourced from a map service. "
            "For landmark_1/2/3 return exactly these strings unchanged:\n"
            + "\n".join(f"  {i+1}. {l}" for i, l in enumerate(real_landmarks))
        )
        # Use json.dumps to safely escape any special characters in landmark strings
        landmark_format = (
            f'"landmark_1": {json.dumps(real_landmarks[0])},\n'
            f'  "landmark_2": {json.dumps(real_landmarks[1])},\n'
            f'  "landmark_3": {json.dumps(real_landmarks[2])}'
        )
    else:
        landmark_instruction = (
            "Real map lookup was unavailable. Use your knowledge of this city to name "
            "3 specific, well-known nearby landmarks within 5km. "
            'Format each as "Landmark Name – 0.Xkm" (max 5km).'
        )
        landmark_format = (
            '"landmark_1": "Landmark Name – 0.Xkm",\n'
            '  "landmark_2": "Landmark Name – 0.Xkm",\n'
            '  "landmark_3": "Landmark Name – 0.Xkm"'
        )

    prompt = f"""You are writing punchy, professional copy for an OOH (Out-of-Home) advertising proposal.

Site details:
- Name: {site_name}
- Location / Address: {location}
- City / Market: {market}
- Format: {fmt}
- Size: {size}

Return ONLY valid JSON (no markdown fences, no extra text) with exactly these keys:

{{
  "tagline": "<4–7 word punchy advertising tagline for this site>",
  "location_desc": "<2–3 sentences describing where the site is and what surrounds it>",
  "visibility_desc": "<2–3 sentences about viewing angles, physical size, and sightlines>",
  "audience_desc": "<2–3 sentences about who passes by and approximate daily volume>",
  {landmark_format}
}}

{landmark_instruction}"""

    response = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=800,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = response.content[0].text.strip()
    raw = re.sub(r"^```[a-z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw)
    text = raw.strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        # Extract the outermost {...} block and retry
        m = re.search(r'\{[\s\S]*\}', text)
        if m:
            try:
                return json.loads(m.group())
            except json.JSONDecodeError:
                pass
        raise ValueError(f"AI returned invalid JSON: {text[:200]}")


# ---------------------------------------------------------------------------
# Helper: build replacement maps from a site plan dict
# ---------------------------------------------------------------------------

def _build_replacements(site: dict) -> tuple[dict, dict]:
    """Return (replacements, ordered) dicts for replace_text_in_slide."""
    lm = [site.get("landmark_1", ""), site.get("landmark_2", ""), site.get("landmark_3", "")]

    replacements = {
        # xyz-format template tokens
        "Site Name":       site.get("site_name", ""),
        "Headline":        site.get("tagline", ""),
        "Size: xyz":       f"Size: {site.get('size', '')}",
        "Format: xyz":     f"Format: {site.get('format', '')}",
        "Location: xyz":   f"Location: {site.get('location', '')}",
        "Frequency: xyz":  f"Frequency: {site.get('frequency', '')}",
        "Units: xyz":      f"Units: {site.get('units', '')}",
        "Traffic: xyz":    f"Traffic: {site.get('traffic', '')}",
        # {TOKEN} style
        "{SITE_NAME}":      site.get("site_name", ""),
        "{TAGLINE}":        site.get("tagline", ""),
        "{LOCATION_DESC}":  site.get("location_desc", ""),
        "{VISIBILITY_DESC}": site.get("visibility_desc", ""),
        "{AUDIENCE_DESC}":  site.get("audience_desc", ""),
        "{SIZE}":           site.get("size", ""),
        "{LOCATION}":       site.get("location", ""),
        "{UNITS}":          site.get("units", ""),
        "{FORMAT}":         site.get("format", ""),
        "{FREQUENCY}":      site.get("frequency", ""),
        "{TRAFFIC}":        site.get("traffic", ""),
        "{LANDMARK_1}":     lm[0],
        "{LANDMARK_2}":     lm[1],
        "{LANDMARK_3}":     lm[2],
        "{MARKET}":         site.get("market", ""),
    }

    ordered = {
        "Text": [
            site.get("location_desc", ""),
            site.get("visibility_desc", ""),
            site.get("audience_desc", ""),
        ],
        "Xyz - 0.5km":     lm,
        "Xyz – 0.5km": lm,
        "Xyz –0.5km":  lm,
        "Xyz -0.5km":      lm,
    }

    return replacements, ordered


# ---------------------------------------------------------------------------
# PPTX build worker (shared by both plan-based and legacy one-shot flows)
# ---------------------------------------------------------------------------

def build_pptx_from_plan(job_id: str, pptx_path: Path, plan: list):
    """Background job: build PPTX from a pre-computed plan list."""
    def update(status: str, message: str, progress: int = 0):
        with jobs_lock:
            jobs[job_id]["status"] = status
            jobs[job_id]["message"] = message
            jobs[job_id]["progress"] = progress

    try:
        update("building", "Loading template…", 5)
        prs = Presentation(str(pptx_path))
        if not prs.slides:
            raise ValueError("The PowerPoint template has no slides.")

        template_slide  = prs.slides[0]
        template_layout = template_slide.slide_layout
        template_spTree = copy.deepcopy(template_slide.shapes._spTree)
        template_spTree_xml = etree.tostring(template_spTree, encoding="unicode")

        # Grab server-side map paths (not sent to client)
        with jobs_lock:
            map_paths = dict(jobs.get(job_id, {}).get("map_paths", {}))

        total = len(plan)

        for idx, site in enumerate(plan):
            pct = 5 + int((idx / total) * 90)
            update("building", f"Building slide {idx + 1}/{total}: {site.get('site_name', '')}…", pct)

            replacements, ordered = _build_replacements(site)

            if idx == 0:
                slide = prs.slides[0]
            else:
                slide = prs.slides.add_slide(template_layout)

                rId_map = {}
                for rel_id, rel in template_slide.part.rels.items():
                    if "image" in rel.reltype or "media" in rel.reltype:
                        try:
                            new_rId = slide.part.relate_to(rel.target_part, rel.reltype)
                            if new_rId != rel_id:
                                rId_map[rel_id] = new_rId
                        except Exception:
                            pass

                xml = template_spTree_xml
                for old_id, new_id in rId_map.items():
                    xml = xml.replace(f'r:embed="{old_id}"', f'r:embed="{new_id}"')
                    xml = xml.replace(f'r:link="{old_id}"',  f'r:link="{new_id}"')

                new_tree = slide.shapes._spTree
                for child in list(new_tree):
                    new_tree.remove(child)
                for child in etree.fromstring(xml):
                    new_tree.append(copy.deepcopy(child))

            replace_text_in_slide(slide, replacements, ordered)

            # Add Google Maps screenshot if available
            map_path_str = map_paths.get(idx)
            if map_path_str:
                map_file = Path(map_path_str)
                if map_file.exists():
                    try:
                        slide.shapes.add_picture(
                            str(map_file),
                            MAP_IMG_LEFT, MAP_IMG_TOP,
                            MAP_IMG_WIDTH, MAP_IMG_HEIGHT,
                        )
                        print(f"[PPTX] map added to slide {idx}")
                    except Exception as pic_err:
                        print(f"[WARN] add_picture failed for slide {idx}: {pic_err}")
                else:
                    print(f"[WARN] map file missing on disk for slide {idx}: {map_path_str}")
            else:
                print(f"[PPTX] no map path for slide {idx} — skipping")

        update("building", "Saving output file…", 96)
        output_filename = f"OOH_Proposal_{job_id[:8]}.pptx"
        output_path = OUTPUT_FOLDER / output_filename
        prs.save(str(output_path))

        with jobs_lock:
            jobs[job_id]["status"]   = "done"
            jobs[job_id]["message"]  = f"Done! {total} slide(s) generated."
            jobs[job_id]["progress"] = 100
            jobs[job_id]["output"]   = output_filename

    except Exception as exc:
        with jobs_lock:
            jobs[job_id]["status"]   = "error"
            jobs[job_id]["message"]  = f"Error: {exc}"
            jobs[job_id]["progress"] = 0
        print(traceback.format_exc())

    finally:
        try:
            pptx_path.unlink(missing_ok=True)
        except Exception:
            pass
        _cleanup_map_images(job_id)


# ---------------------------------------------------------------------------
# Plan generation worker
# ---------------------------------------------------------------------------

def _ensure_site_name_column(df: pd.DataFrame) -> pd.DataFrame:
    """Some site plans (e.g. Etihad's) use 'Location' as both the site name and
    the location description. If there's no 'Site Name' column but there is a
    'Location' column, treat Location's value as the site name too."""
    if "Site Name" not in df.columns and "Location" in df.columns:
        df["Site Name"] = df["Location"]
    return df


_EXCEL_HEADER_KEYWORDS = {
    'site name', 'location', 'market', 'format', 'size',
    'units/faces', 'spot duration', 'sov/loop', 'impacts',
}


def _detect_excel_header_row(excel_path) -> int:
    """Find which row holds the real column headers, scanning the first 25 rows.

    Handles plans with a title/client/campaign block above the header (like
    Etihad's) by requiring at least 2 recognized column names. Falls back to
    row 0 — pandas' normal default — when nothing better is found, so plain
    flat-header files behave exactly as before.
    """
    import openpyxl as opx
    wb = opx.load_workbook(excel_path, data_only=True, read_only=True)
    ws = wb.active
    for i, row in enumerate(ws.iter_rows(max_row=25, values_only=True)):
        vals = {str(v).strip().lower() for v in row if v is not None and str(v).strip()}
        if len(vals & _EXCEL_HEADER_KEYWORDS) >= 2:
            wb.close()
            return i
    wb.close()
    return 0


def generate_plan_job(job_id: str, excel_path: Path):
    """Background job: read Excel + AI/landmarks, produce a content plan."""
    def update(status: str, message: str, progress: int = 0):
        with jobs_lock:
            jobs[job_id]["status"]   = status
            jobs[job_id]["message"]  = message
            jobs[job_id]["progress"] = progress

    try:
        update("planning", "Reading Excel file…", 5)
        header_row = _detect_excel_header_row(excel_path)
        df = pd.read_excel(excel_path, engine="openpyxl", header=header_row)
        df.columns = [c.strip() for c in df.columns]

        if "Market" in df.columns:
            df["Market"] = df["Market"].ffill()

        df = _ensure_site_name_column(df)
        if "Site Name" not in df.columns:
            raise ValueError("Excel file must have a 'Site Name' column.")

        df = df[df["Site Name"].notna() & (df["Site Name"].astype(str).str.strip() != "")]
        df = df.reset_index(drop=True)

        if df.empty:
            raise ValueError("No valid site rows found in the Excel file.")

        total = len(df)
        update("planning", f"Found {total} site(s). Connecting to AI…", 10)

        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY environment variable is not set.")
        client = anthropic.Anthropic(api_key=api_key)

        plan = []
        for idx, row in df.iterrows():
            pct = 10 + int(((idx + 1) / total) * 85)
            site_name = str(row.get("Site Name", "")).strip()
            update("planning", f"Researching site {idx + 1}/{total}: {site_name}…", pct)

            location  = str(row.get("Location", "")).strip()
            market    = str(row.get("Market", "")).strip()
            is_mobile = location.lower() == "various"

            spot_dur = str(row.get("Spot Duration", "")).strip()
            sov_loop = str(row.get("SOV/Loop", "")).strip()
            if spot_dur.lower() in ("", "nan", "n/a", "na"):
                frequency = sov_loop
            else:
                frequency = f"{spot_dur} {sov_loop}".strip()

            raw_impacts = row.get("Impacts", "")
            try:
                traffic = f"{int(float(str(raw_impacts).replace(',', ''))):,}"
            except (ValueError, TypeError):
                traffic = str(raw_impacts).strip()

            try:
                ai = generate_site_content(row.to_dict(), client)
            except Exception as ai_err:
                print(f"[WARN] AI failed for site {site_name!r}: {ai_err}")
                ai = {
                    "tagline":         "",
                    "location_desc":   "",
                    "visibility_desc": "",
                    "audience_desc":   "",
                    "landmark_1":      "",
                    "landmark_2":      "",
                    "landmark_3":      "",
                }

            # Fetch map screenshot
            map_address = market if is_mobile else location
            map_zoom    = 11 if is_mobile else 16
            print(f"[MAP] fetching for site {idx} '{site_name}': address={map_address!r}, city={market!r}, zoom={map_zoom}")
            map_bytes   = get_map_image_bytes(map_address, market, zoom=map_zoom)
            has_map     = False
            if map_bytes:
                map_file = UPLOAD_FOLDER / f"{job_id}_map_{idx}.png"
                map_file.write_bytes(map_bytes)
                with jobs_lock:
                    jobs[job_id]["map_paths"][idx] = str(map_file)
                has_map = True
                print(f"[MAP] saved {len(map_bytes)} bytes for site {idx} '{site_name}'")
            else:
                print(f"[MAP] no map returned for site {idx} '{site_name}'")

            plan.append({
                "site_name":       site_name,
                "market":          market,
                "location":        location,
                "format":          str(row.get("Format", "")).strip(),
                "size":            str(row.get("Size", "")).strip(),
                "units":           str(row.get("Units/Faces", "")).strip(),
                "frequency":       frequency,
                "traffic":         traffic,
                "tagline":         ai.get("tagline", ""),
                "location_desc":   ai.get("location_desc", ""),
                "visibility_desc": ai.get("visibility_desc", ""),
                "audience_desc":   ai.get("audience_desc", ""),
                "landmark_1":      ai.get("landmark_1", ""),
                "landmark_2":      ai.get("landmark_2", ""),
                "landmark_3":      ai.get("landmark_3", ""),
                "has_map":         has_map,
            })

        with jobs_lock:
            jobs[job_id]["status"]   = "plan_ready"
            jobs[job_id]["message"]  = f"Plan ready — {total} site(s). Review and edit, then build."
            jobs[job_id]["progress"] = 100
            jobs[job_id]["plan"]     = plan

    except Exception as exc:
        with jobs_lock:
            jobs[job_id]["status"]   = "error"
            jobs[job_id]["message"]  = f"Error: {exc}"
            jobs[job_id]["progress"] = 0
        print(traceback.format_exc())

    finally:
        try:
            excel_path.unlink(missing_ok=True)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Legacy one-shot worker (kept for /api/generate backward compat)
# ---------------------------------------------------------------------------

def process_job(job_id: str, excel_path: Path, pptx_path: Path):
    """Upload + AI + build in one shot (legacy endpoint)."""
    def update(status: str, message: str, progress: int = 0):
        with jobs_lock:
            jobs[job_id]["status"]   = status
            jobs[job_id]["message"]  = message
            jobs[job_id]["progress"] = progress

    try:
        update("processing", "Reading Excel file…", 5)
        header_row = _detect_excel_header_row(excel_path)
        df = pd.read_excel(excel_path, engine="openpyxl", header=header_row)
        df.columns = [c.strip() for c in df.columns]

        if "Market" in df.columns:
            df["Market"] = df["Market"].ffill()

        df = _ensure_site_name_column(df)
        if "Site Name" not in df.columns:
            raise ValueError("Excel file must have a 'Site Name' column.")
        df = df[df["Site Name"].notna() & (df["Site Name"].astype(str).str.strip() != "")]
        df = df.reset_index(drop=True)

        if df.empty:
            raise ValueError("No valid site rows found in the Excel file.")

        total = len(df)
        update("processing", f"Found {total} site(s). Loading template…", 10)

        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY environment variable is not set.")
        client = anthropic.Anthropic(api_key=api_key)

        plan = []
        for idx, row in df.iterrows():
            pct = 10 + int(((idx + 1) / total) * 80)
            site_name = str(row.get("Site Name", "")).strip()
            update("processing", f"Processing {idx + 1}/{total}: {site_name}…", pct)

            spot_dur = str(row.get("Spot Duration", "")).strip()
            sov_loop = str(row.get("SOV/Loop", "")).strip()
            frequency = sov_loop if spot_dur.lower() in ("", "nan", "n/a", "na") else f"{spot_dur} {sov_loop}".strip()

            raw_impacts = row.get("Impacts", "")
            try:
                traffic = f"{int(float(str(raw_impacts).replace(',', ''))):,}"
            except (ValueError, TypeError):
                traffic = str(raw_impacts).strip()

            ai = generate_site_content(row.to_dict(), client)

            plan.append({
                "site_name":       site_name,
                "market":          str(row.get("Market", "")).strip(),
                "location":        str(row.get("Location", "")).strip(),
                "format":          str(row.get("Format", "")).strip(),
                "size":            str(row.get("Size", "")).strip(),
                "units":           str(row.get("Units/Faces", "")).strip(),
                "frequency":       frequency,
                "traffic":         traffic,
                "tagline":         ai.get("tagline", ""),
                "location_desc":   ai.get("location_desc", ""),
                "visibility_desc": ai.get("visibility_desc", ""),
                "audience_desc":   ai.get("audience_desc", ""),
                "landmark_1":      ai.get("landmark_1", ""),
                "landmark_2":      ai.get("landmark_2", ""),
                "landmark_3":      ai.get("landmark_3", ""),
            })

        excel_path.unlink(missing_ok=True)
        build_pptx_from_plan(job_id, pptx_path, plan)

    except Exception as exc:
        with jobs_lock:
            jobs[job_id]["status"]   = "error"
            jobs[job_id]["message"]  = f"Error: {exc}"
            jobs[job_id]["progress"] = 0
        print(traceback.format_exc())
        try:
            excel_path.unlink(missing_ok=True)
            pptx_path.unlink(missing_ok=True)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.route("/")
def index():
    return render_template("index.html")


# ── Pro flow ────────────────────────────────────────────────────────────────

@app.route("/api/plan", methods=["POST"])
def create_plan():
    """Step 1: Upload files, generate content plan (AI + landmarks)."""
    if "excel" not in request.files or "template" not in request.files:
        return jsonify({"error": "Both 'excel' and 'template' files are required."}), 400

    excel_file    = request.files["excel"]
    template_file = request.files["template"]

    if not excel_file.filename.endswith((".xlsx", ".xls")):
        return jsonify({"error": "Excel file must be .xlsx or .xls"}), 400
    if not template_file.filename.endswith(".pptx"):
        return jsonify({"error": "Template file must be .pptx"}), 400

    job_id     = uuid.uuid4().hex
    excel_path = UPLOAD_FOLDER / f"{job_id}_data.xlsx"
    pptx_path  = UPLOAD_FOLDER / f"{job_id}_template.pptx"
    excel_file.save(str(excel_path))
    template_file.save(str(pptx_path))

    with jobs_lock:
        jobs[job_id] = {
            "status":    "planning",
            "message":   "Starting…",
            "progress":  0,
            "plan":      None,
            "map_paths": {},        # {site_idx: str(path)} — server-side only
            "pptx_path": str(pptx_path),
            "output":    None,
        }

    threading.Thread(target=generate_plan_job, args=(job_id, excel_path), daemon=True).start()
    return jsonify({"job_id": job_id})


@app.route("/api/build", methods=["POST"])
def build():
    """Step 2: Submit (possibly edited) plan → build PPTX."""
    data = request.get_json(force=True, silent=True) or {}
    job_id = data.get("job_id")
    plan   = data.get("plan")

    if not job_id or not plan:
        return jsonify({"error": "job_id and plan are required"}), 400

    with jobs_lock:
        job = jobs.get(job_id)

    if not job:
        return jsonify({"error": "Job not found — session may have expired."}), 404

    pptx_path = Path(job.get("pptx_path", ""))
    if not pptx_path.exists():
        return jsonify({"error": "Template file not found. Please re-upload and start again."}), 404

    with jobs_lock:
        jobs[job_id]["status"]   = "building"
        jobs[job_id]["message"]  = "Starting build…"
        jobs[job_id]["progress"] = 0
        jobs[job_id]["output"]   = None

    threading.Thread(target=build_pptx_from_plan, args=(job_id, pptx_path, plan), daemon=True).start()
    return jsonify({"job_id": job_id})


# ── Shared polling + download ────────────────────────────────────────────────

@app.route("/api/status/<job_id>")
def status(job_id: str):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return jsonify({"error": "Job not found"}), 404
    # Only return fields safe for the frontend (omit pptx_path)
    return jsonify({
        "status":   job["status"],
        "message":  job["message"],
        "progress": job["progress"],
        "plan":     job.get("plan"),
        "output":   job.get("output"),
    })


@app.route("/api/download/<job_id>")
def download(job_id: str):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job or job["status"] != "done":
        return jsonify({"error": "File not ready"}), 404

    output_path = OUTPUT_FOLDER / job["output"]
    if not output_path.exists():
        return jsonify({"error": "Output file missing"}), 404

    response = send_file(
        str(output_path),
        as_attachment=True,
        download_name=job["output"],
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    @response.call_on_close
    def _cleanup():
        try:
            output_path.unlink(missing_ok=True)
        except Exception:
            pass
        with jobs_lock:
            jobs.pop(job_id, None)

    return response


# ── Map image preview (served to the plan review UI) ────────────────────────

@app.route("/api/map/<job_id>/<int:site_idx>")
def serve_map(job_id: str, site_idx: int):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return "", 404
    map_path_str = job.get("map_paths", {}).get(site_idx)
    if not map_path_str:
        return "", 404
    p = Path(map_path_str)
    if not p.exists():
        return "", 404
    return send_file(str(p), mimetype="image/png")


# ── Legacy one-shot endpoint (backward compat) ───────────────────────────────

@app.route("/api/generate", methods=["POST"])
def generate():
    if "excel" not in request.files or "template" not in request.files:
        return jsonify({"error": "Both 'excel' and 'template' files are required."}), 400

    excel_file    = request.files["excel"]
    template_file = request.files["template"]

    if not excel_file.filename.endswith((".xlsx", ".xls")):
        return jsonify({"error": "Excel file must be .xlsx or .xls"}), 400
    if not template_file.filename.endswith(".pptx"):
        return jsonify({"error": "Template file must be .pptx"}), 400

    job_id     = uuid.uuid4().hex
    excel_path = UPLOAD_FOLDER / f"{job_id}_data.xlsx"
    pptx_path  = UPLOAD_FOLDER / f"{job_id}_template.pptx"
    excel_file.save(str(excel_path))
    template_file.save(str(pptx_path))

    with jobs_lock:
        jobs[job_id] = {
            "status":    "queued",
            "message":   "Queued…",
            "progress":  0,
            "plan":      None,
            "pptx_path": str(pptx_path),
            "output":    None,
        }

    threading.Thread(target=process_job, args=(job_id, excel_path, pptx_path), daemon=True).start()
    return jsonify({"job_id": job_id})


# ---------------------------------------------------------------------------
# CN Print Plan Filler
# ---------------------------------------------------------------------------
# Fills the bundled Conde Nast print plan template (assets/) from up to three
# publisher rate cards (India / UK / US). All extraction and template surgery
# lives in cn_print_plan.py; see that module for the section geometry.

from cn_print_plan import fill_cn_print_plan, CNPlanError

CN_TEMPLATE_PATH = BASE_DIR / "assets" / "CN_Print_Plan_Template_EMPTY.xlsx"


@app.route('/print-plan')
def print_plan_page():
    return render_template('print_plan.html')


@app.route('/fill-cn-plan', methods=['POST'])
def fill_cn_plan():
    cards = {}
    for mkt in ('india', 'uk', 'us'):
        f = request.files.get(mkt)
        if f and f.filename:
            cards[mkt] = f.read()
    if not cards:
        return 'Upload at least one rate card (India, UK or US)', 400

    tpl = request.files.get('template')
    template_source = io.BytesIO(tpl.read()) if tpl and tpl.filename else CN_TEMPLATE_PATH
    header = {k: (request.form.get(k) or '').strip()
              for k in ('client', 'campaign', 'agency', 'contact')}

    try:
        out = fill_cn_print_plan(template_source, cards, header)
    except CNPlanError as e:
        return str(e), 400
    except Exception:
        print(traceback.format_exc())
        return 'Failed to fill the plan — check the server logs', 500

    campaign = re.sub(r'[^A-Za-z0-9 _\-]', '', header.get('campaign', '')).strip()
    fname = (campaign.replace(' ', '_') + '_Print_Plan.xlsx') if campaign else 'CN_Print_Plan_Filled.xlsx'
    return send_file(out, as_attachment=True, download_name=fname,
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')


# ---------------------------------------------------------------------------
# Unified OOH Proposal Deck Builder
# ---------------------------------------------------------------------------
# Fills a branded 2-slide template (cover + one site page) from a vendor
# Excel media plan (one row per site: Market/Site Name/Location/Format/
# Units/Size/Spot Duration/SOV-Loop/Impacts). The site-page slide is cloned
# once per site; there's no vendor photo, so its "[SITE PHOTO]"/"[MAP]"
# placeholders are simply dropped. The slide title always matches the
# plan's Site Name verbatim. AI fills the descriptive copy (Location/
# Visibility/Audience/Why this site) and real landmark lookups (with an
# AI fallback when geocoding finds nothing) run per site — network-bound
# work, so the build runs as a background job polled via the existing
# /api/status and /api/download endpoints rather than a single
# synchronous request.

_OOH_NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_OOH_NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

# Header text -> internal field name. Matched case-insensitively against the
# Excel's header row, wherever that row happens to land (there's a Date/
# Client/Campaign block above it).
_OOH_HEADER_KEYWORDS = {
    'market':            ['market'],
    'site_name':         ['site name'],
    'location':          ['location'],
    'format':            ['format'],
    'units':             ['units'],
    'size':              ['size'],
    'spot_duration':     ['spot duration'],
    'sov_loop':          ['sov/loop', 'sov / loop'],
    'campaign_duration': ['campaign duration'],
    'impacts':           ['impacts'],
}
# Market is blank on continuation rows within the same city section — every
# other column (including Site Name and Location) is filled on every row.
_OOH_FORWARD_FILL_FIELDS = ('market',)


def _ooh_find_header_row(rows: list) -> int:
    for i, row in enumerate(rows[:20]):
        vals = {str(v).strip().lower() for v in row if v is not None and str(v).strip()}
        if 'location' in vals and 'market' in vals:
            return i
    raise ValueError("Could not find the header row (expected 'Location' and 'Market' columns) in the Excel file.")


def _ooh_parse_excel_sites(excel_bytes: bytes) -> list:
    """Parse the vendor media-plan Excel into a flat list of site dicts, one per data row."""
    import openpyxl as opx
    wb = opx.load_workbook(io.BytesIO(excel_bytes), data_only=True)
    ws = wb.active
    raw = [[c.value for c in row] for row in ws.iter_rows()]
    if not raw:
        raise ValueError("The Excel file is empty.")

    hdr_idx = _ooh_find_header_row(raw)
    headers = [str(h).strip().lower() if h is not None else '' for h in raw[hdr_idx]]

    field_col = {}
    for field, keywords in _OOH_HEADER_KEYWORDS.items():
        for ci, h in enumerate(headers):
            if h in keywords:
                field_col[field] = ci
                break

    if 'location' not in field_col or 'site_name' not in field_col:
        raise ValueError("Could not find 'Site Name' and 'Location' columns in the Excel file.")

    def cell(row, field):
        ci = field_col.get(field)
        if ci is None or ci >= len(row) or row[ci] is None:
            return ''
        return str(row[ci]).strip()

    sites = []
    last = {f: '' for f in _OOH_FORWARD_FILL_FIELDS}
    sno_counter = 0

    for row in raw[hdr_idx + 1:]:
        if not cell(row, 'site_name'):
            continue  # blank separator row, subtotal, or trailing terms/disclaimer text — none of those have a Site Name

        values = {}
        for field in _OOH_FORWARD_FILL_FIELDS:
            v = cell(row, field)
            values[field] = v or last[field]
            last[field] = values[field]

        sno_counter += 1
        sites.append({
            'sno':               sno_counter,
            'market':            values['market'] or 'Other',
            'site_name':         cell(row, 'site_name'),
            'location':          cell(row, 'location'),
            'format':            cell(row, 'format'),
            'units':             cell(row, 'units') or '1',
            'size':              cell(row, 'size'),
            'spot_duration':     cell(row, 'spot_duration'),
            'sov_loop':          cell(row, 'sov_loop'),
            'campaign_duration': cell(row, 'campaign_duration'),
            'impacts':           cell(row, 'impacts'),
        })

    if not sites:
        raise ValueError("No site rows found under the header row.")
    return sites


def _ooh_format_site_fields(site: dict) -> dict:
    """Derive the plain-fact display strings (Format/Size/Units/etc.) from raw Excel fields.

    Format/Size/Spot Duration/SOV-Loop already arrive as ready-to-display free
    text from this sheet (e.g. "54,4m x 13,44m", "60 spots per day") — no
    parsing needed, unlike the old column set this replaced.
    """
    impacts_digits = re.sub(r'[^\d.]', '', site['impacts'])
    try:
        impacts_fmt = f"{int(float(impacts_digits)):,}" if impacts_digits else site['impacts']
    except ValueError:
        impacts_fmt = site['impacts']

    visibility_fallback = ', '.join(p for p in (site['format'], site['size'], site['spot_duration']) if p)

    return {
        'format':        site['format'],
        'size':          site['size'],
        'units':         site['units'],
        'spot_duration': site['spot_duration'],
        'sov':           site['sov_loop'],
        'traffic':       impacts_fmt,
        # Fallbacks used only if AI content is unavailable (no API key, or the call failed).
        # Kept short since this fills a large single-line title font (a longer
        # string wraps and overlaps the subtitle below it).
        'site_name_fallback':   site['site_name'][:24],
        'location_fallback':    site['location'],
        'visibility_fallback':  visibility_fallback,
        'audience_fallback':    impacts_fmt,
    }


def _ooh_generate_ai_content(site: dict, client_name: str, ai_client: anthropic.Anthropic) -> dict:
    """Ask Claude for the site's display name/nickname, 2-3 sentence descriptive
    copy for Location/Visibility/Audience, and a one-line "why this site" pitch
    tailored to the client."""
    prompt = f"""You are writing punchy, professional copy for an OOH (Out-of-Home) advertising strategy proposal for the client "{client_name}".

Site details:
- Site name: {site['site_name']}
- Location: {site['location']}
- Market/City: {site['market']}
- Format: {site['format']}
- Size: {site['size']}
- Spot duration: {site['spot_duration']}
- SOV/Loop: {site['sov_loop']}
- Units: {site['units']}
- Impacts: {site['impacts']}

The site's title on the slide is always "{site['site_name']}" verbatim from the plan — do not rename or paraphrase it.

Return ONLY valid JSON (no markdown fences, no extra text) with exactly these keys:

{{
  "site_nickname": "<a short, catchy 2-5 word descriptor/nickname for this specific site, e.g. 'The Ball Drop Tower'>",
  "location_desc": "<2-3 sentences describing where this site is located and its surroundings>",
  "visibility_desc": "<2-3 sentences about the screen's visibility, format, and viewing conditions>",
  "audience_desc": "<2-3 sentences about the audience/traffic this site reaches>",
  "why_this_site": "<1-2 sentences on why this specific site is a strong fit for {client_name}>"
}}"""

    response = ai_client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=600,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text.strip()
    raw = re.sub(r"^```[a-z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw)
    text = raw.strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        m = re.search(r'\{[\s\S]*\}', text)
        if m:
            return json.loads(m.group())
        raise ValueError(f"AI returned invalid JSON: {text[:200]}")


def _ooh_generate_ai_landmarks(location: str, market: str, ai_client: anthropic.Anthropic, n: int = 3) -> list:
    """Fallback for when real geocoding can't find nearby landmarks (the free
    OSM/Overpass lookup is unreliable — it requires >= n hits within a tight
    radius and can fail even for a whole city name, not just a noisy address).
    Asks Claude to name real, well-known landmarks from its own knowledge,
    in the same "Name - 0.Xkm" format the real lookup produces."""
    prompt = f"""Name {n} real, well-known landmarks near this location, using your general knowledge of the area.

Location: {location}
City: {market}

Return ONLY valid JSON (no markdown fences, no extra text), a single array of exactly {n} strings, each formatted as "Landmark Name - 0.Xkm" with a plausible walking distance estimate (under 2km). Example: ["Rotterdam Bridge - 0.5km", "Central Station - 1.2km", "City Park - 1.8km"]"""

    response = ai_client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=300,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text.strip()
    raw = re.sub(r"^```[a-z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw)
    text = raw.strip()
    try:
        result = json.loads(text)
    except json.JSONDecodeError:
        m = re.search(r'\[[\s\S]*\]', text)
        if not m:
            raise ValueError(f"AI returned invalid JSON: {text[:200]}")
        result = json.loads(m.group())
    return [str(x) for x in result][:n]


def _ooh_clone_slide(prs: Presentation, idx: int):
    """Clone slide idx, preserving any embedded images (e.g. the Skyscale logo)
    by re-relating each image part to the new slide and remapping r:embed/
    r:link IDs in the copied shape XML — a freshly-created relationship isn't
    guaranteed to get the same ID as the original, so a naive deep-copy of
    the shape tree leaves dangling references that show up as broken/missing
    images in PowerPoint even though the file opens without error."""
    tpl = prs.slides[idx]
    new = prs.slides.add_slide(tpl.slide_layout)

    rId_map = {}
    for rel_id, rel in tpl.part.rels.items():
        if 'image' in rel.reltype:
            try:
                new_rId = new.part.relate_to(rel.target_part, rel.reltype)
                if new_rId != rel_id:
                    rId_map[rel_id] = new_rId
            except Exception:
                pass

    xml = etree.tostring(tpl.shapes._spTree, encoding='unicode')
    for old_id, new_id in rId_map.items():
        xml = xml.replace(f'r:embed="{old_id}"', f'r:embed="{new_id}"')
        xml = xml.replace(f'r:link="{old_id}"', f'r:link="{new_id}"')

    new_tree = new.shapes._spTree
    for child in list(new_tree):
        new_tree.remove(child)
    for child in etree.fromstring(xml):
        new_tree.append(copy.deepcopy(child))

    return new


def _ooh_shape_text(shape) -> str:
    try:
        return shape.text_frame.text.strip()
    except Exception:
        return ''


def _ooh_set_shape_text(shape, text: str):
    """Replace a shape's visible text, preserving its first run's existing
    formatting (font/size/color) instead of rebuilding a run from scratch —
    this template's shapes are one label/value per shape with real design
    styling already applied, unlike the multi-paragraph Etihad template."""
    tf = shape.text_frame
    paras = tf.paragraphs
    if not paras:
        return
    p0 = paras[0]
    if p0.runs:
        p0.runs[0].text = text
        for extra_run in p0.runs[1:]:
            extra_run.text = ''
    else:
        p0.text = text
    # Drop any extra (typically empty) trailing paragraphs so nothing stray remains.
    for extra_p in list(paras[1:]):
        el = extra_p._p
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


def _ooh_has_blip_fill(shape) -> bool:
    """True if this shape's fill is a picture (blipFill) — used to find the
    photo/map placeholders, which are custom-geometry shapes with an image
    fill rather than plain <p:pic> elements (typical of a Canva export)."""
    try:
        spPr = shape._element.spPr
    except Exception:
        return False
    if spPr is None:
        return False
    return spPr.find(qn('a:blipFill')) is not None


_OOH_STATIC_LABELS = {'LOCATION', 'VISIBILITY', 'AUDIENCE', 'TECHNICAL SPECIFICATIONS', 'NEARBY LANDMARKS'}
_OOH_PREFIXED_FIELDS = [
    ('FORMAT :', 'format'),
    ('SIZE :', 'size'),
    ('UNITS :', 'units'),
    ('SPOT LENGTH :', 'spot_duration'),
    ('FREQUENCY :', 'sov'),
    ('TRAFFIC :', 'traffic'),
]


def _ooh_fill_cover_slide(slide, client_name: str, campaign_subtitle: str, campaign_name: str, duration: str):
    """Fill the cover slide's 4 text shapes, identified by sorting them by
    vertical position (client name, subtitle, campaign name, duration — in
    that reading order, verified against the actual template file)."""
    candidates = [sh for sh in slide.shapes if _ooh_shape_text(sh)]
    candidates.sort(key=lambda sh: sh.top)
    values = [client_name, campaign_subtitle, campaign_name, duration]
    for sh, value in zip(candidates, values):
        _ooh_set_shape_text(sh, value)


def _ooh_fill_site_slide(slide, fields: dict, ai: dict, client_name: str, landmarks: list):
    """Fill a cloned site-page slide.

    Labels (LOCATION/VISIBILITY/AUDIENCE/etc.) and prefixed fact lines
    (FORMAT :/SIZE :/etc.) are matched by their own stable text. The title,
    subtitle, four AI-prose values (why-this-site, location, visibility,
    audience) and three landmark lines have no stable per-field marker of
    their own — matched instead by sorting them by position, which reliably
    reproduces the template's label/value layout (verified against the
    actual template file's real coordinates).
    """
    placeholder_shapes = []  # [SITE PHOTO] and [MAP] boxes — no image source for either, so just removed
    left_column = []   # title + subtitle candidates
    leftover = []      # why-this-site / location / visibility / audience / 3 landmarks

    for sh in slide.shapes:
        if _ooh_has_blip_fill(sh):
            if sh.width > 8_000_000 and sh.left < 1_000_000:
                placeholder_shapes.append(sh)
            continue

        text = _ooh_shape_text(sh)
        if not text:
            continue
        upper = text.upper()

        if upper in _OOH_STATIC_LABELS:
            continue  # static label, no per-site change needed

        if upper.startswith('WHY THIS SITE FOR'):
            _ooh_set_shape_text(sh, f"Why this Site for {client_name}")
            continue

        matched = False
        for prefix, field in _OOH_PREFIXED_FIELDS:
            if upper.startswith(prefix):
                _ooh_set_shape_text(sh, f"{prefix} {fields.get(field, '')}".strip())
                matched = True
                break
        if matched:
            continue

        if sh.left < 1_000_000:
            left_column.append(sh)
        else:
            leftover.append(sh)

    left_column.sort(key=lambda sh: sh.top)
    if len(left_column) >= 1:
        # Always the raw Excel Site Name (not an AI paraphrase) so the slide
        # title matches the plan exactly. Capped since this is a large
        # single-line title font — a longer string wraps and overlaps the
        # subtitle sitting right below it.
        title = fields.get('site_name_fallback', '')[:24]
        _ooh_set_shape_text(left_column[0], title)
    if len(left_column) >= 2:
        nickname = ai.get('site_nickname', '')
        market = fields.get('market', '')
        _ooh_set_shape_text(left_column[1], f"{market.upper()}  |  {nickname}" if nickname else market.upper())

    leftover.sort(key=lambda sh: sh.top)
    prose_values = [
        ai.get('why_this_site') or f"A strong fit for {client_name}'s target audience.",
        ai.get('location_desc') or fields.get('location_fallback', ''),
        ai.get('visibility_desc') or fields.get('visibility_fallback', ''),
        ai.get('audience_desc') or fields.get('audience_fallback', ''),
    ]
    for sh, value in zip(leftover[:4], prose_values):
        _ooh_set_shape_text(sh, value)
    for sh, value in zip(leftover[4:7], (landmarks + ['', '', ''])[:3]):
        _ooh_set_shape_text(sh, value)

    spt = slide.shapes._spTree
    for sh in placeholder_shapes:
        spt.remove(sh._element)  # no vendor photo, and the maps feature is unused — leave blank


def _ooh_build_deck(prs: Presentation, sites: list, client_name: str, campaign_subtitle: str,
                     campaign_name: str, duration: str, per_site: dict):
    """Fill the cover slide (idx 0) in place, clone the site slide (idx 1) once
    per site, then remove the original site-template slide."""
    _ooh_fill_cover_slide(prs.slides[0], client_name, campaign_subtitle, campaign_name, duration)

    for site in sites:
        slide = _ooh_clone_slide(prs, 1)
        data = per_site.get(site['sno'], {})
        _ooh_fill_site_slide(
            slide,
            {**_ooh_format_site_fields(site), 'market': site['market']},
            data.get('ai_content') or {},
            client_name,
            data.get('landmarks') or [],
        )

    # Remove the original site-template slide (idx 1) now that every site has its own clone.
    sll = prs.slides._sldIdLst
    pp = prs.part
    elem = list(sll)[1]
    rId = elem.get(f'{{{_OOH_NS_R}}}id')
    sll.remove(elem)
    if rId:
        try:
            pp._rels.pop(rId)
        except Exception:
            pass


def _ooh_build_job(job_id: str, template_bytes: bytes, sites: list, client_name: str,
                    campaign_subtitle: str, campaign_name: str, duration: str):
    """Background job: AI copy + real landmarks per site, then assemble the deck."""
    def update(status, message, progress=0):
        with jobs_lock:
            jobs[job_id]['status'] = status
            jobs[job_id]['message'] = message
            jobs[job_id]['progress'] = progress

    try:
        total = len(sites)

        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        ai_client = anthropic.Anthropic(api_key=api_key) if api_key else None
        if not ai_client:
            print("[OOH DECK AI] ANTHROPIC_API_KEY not set — skipping AI copy, using raw field values")

        per_site = {}
        for idx, site in enumerate(sites):
            pct = int((idx / total) * 90)
            update('building', f"Researching {idx + 1}/{total}: {site['location'][:60]}…", pct)

            data = {}
            if ai_client:
                try:
                    data['ai_content'] = _ooh_generate_ai_content(site, client_name, ai_client)
                except Exception as e:
                    print(f"[OOH DECK AI] failed for site {site['sno']} ({site['location']!r}): {e}")

            landmarks = []
            try:
                landmarks = get_real_landmarks(site['location'], site['market'], n=3)
                if not landmarks:
                    # The full location text is often too descriptive to
                    # geocode — fall back to city-level landmarks rather than
                    # leaving the section blank.
                    landmarks = get_real_landmarks(site['market'], site['market'], n=3)
            except Exception as e:
                print(f"[OOH DECK LANDMARKS] failed for site {site['sno']} ({site['location']!r}): {e}")

            if not landmarks and ai_client:
                # The free OSM/Overpass lookup is unreliable — it requires >= 3
                # hits within a tight radius and can come back empty even for a
                # whole city name. Fall back to Claude's own knowledge of the area.
                try:
                    landmarks = _ooh_generate_ai_landmarks(site['location'], site['market'], ai_client, n=3)
                except Exception as e:
                    print(f"[OOH DECK LANDMARKS] AI fallback failed for site {site['sno']} ({site['location']!r}): {e}")

            data['landmarks'] = landmarks or []

            per_site[site['sno']] = data

        update('building', 'Assembling slides…', 92)
        prs = Presentation(io.BytesIO(template_bytes))
        _ooh_build_deck(prs, sites, client_name, campaign_subtitle, campaign_name, duration, per_site)

        update('building', 'Saving file…', 97)
        output_filename = f"OOH_Deck_{job_id[:8]}.pptx"
        output_path = OUTPUT_FOLDER / output_filename
        prs.save(str(output_path))

        with jobs_lock:
            jobs[job_id]['status']   = 'done'
            jobs[job_id]['message']  = f"Done! {total} site slide(s) generated."
            jobs[job_id]['progress'] = 100
            jobs[job_id]['output']   = output_filename

    except Exception as exc:
        with jobs_lock:
            jobs[job_id]['status']   = 'error'
            jobs[job_id]['message']  = f"Error: {exc}"
            jobs[job_id]['progress'] = 0
        print(traceback.format_exc())


@app.route('/ooh-deck')
def ooh_deck_page():
    return render_template('ooh_deck.html')


@app.route('/build-ooh-deck', methods=['POST'])
def build_ooh_deck_route():
    """Step 1: upload the Excel site plan + template + client/campaign details, start the build job."""
    excel_file    = request.files.get('excel')
    template_file = request.files.get('template')
    client_name       = (request.form.get('client_name') or '').strip()
    campaign_name     = (request.form.get('campaign_name') or '').strip()
    campaign_subtitle = (request.form.get('campaign_subtitle') or '').strip()
    duration          = (request.form.get('duration') or '').strip()

    if not excel_file or not excel_file.filename:
        return jsonify({'error': 'Excel site-plan file is required.'}), 400
    if not template_file or not template_file.filename:
        return jsonify({'error': 'Template file is required.'}), 400
    if not excel_file.filename.lower().endswith(('.xlsx', '.xls')):
        return jsonify({'error': 'Site plan must be .xlsx or .xls'}), 400
    if not template_file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Template must be a .pptx'}), 400
    if not client_name:
        return jsonify({'error': 'Client name is required.'}), 400

    try:
        sites = _ooh_parse_excel_sites(excel_file.read())
    except Exception as e:
        return jsonify({'error': str(e)}), 400

    template_bytes = template_file.read()
    job_id = uuid.uuid4().hex
    with jobs_lock:
        jobs[job_id] = {
            'status':   'building',
            'message':  f"Found {len(sites)} site(s). Starting…",
            'progress': 0,
            'plan':     None,
            'output':   None,
        }

    threading.Thread(
        target=_ooh_build_job,
        args=(job_id, template_bytes, sites, client_name, campaign_subtitle, campaign_name, duration),
        daemon=True,
    ).start()
    return jsonify({'job_id': job_id, 'site_count': len(sites)})


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    port  = int(os.environ.get("PORT", 5000))
    debug = os.environ.get("FLASK_DEBUG", "0") == "1"
    app.run(host="0.0.0.0", port=port, debug=debug)
